"""
Converts an SVG shape into a Microsoft Office object, and saves as pptx file

"""
import re
from lxml import etree, html
from lxml.builder import ElementMaker
from pypptx import a, p, shape, color, nsmap, cust_shape
from color import rgba


re_ns = re.compile(r'({.*?})?(.*)')
re_path = re.compile(r'[mMzZlLhHvVcCsSqQtTaA]|[\+\-]?[\d\.e]+')

def interpret_str(val):
    if val:
        rn = re.compile(r'([\d\.-]+)')
        match = rn.search(val)
        value = match.group(1)
        if val.startswith('-'):
            val = str(float(value) - float(value))
        elif val.endswith('%'):
            val = str(int(value) * 0.2)
        elif val.endswith('em'):
            val = str(value * 10 + 8)
        elif val.endswith('pt'):
            val = str(int(value + 6))
        else:
            val = value
    return val

def msclr(color):
    r, g, b, a = rgba(color)
    return '%02x%02x%02x' % (255*r, 255*g, 255*b)

def css_style(style):
    e = {}
    attrs = [x for x in style.split(";") if x != '']

    for attr in attrs:
        keys, values = attr.split(':')
        key, value = keys.split(), values.split()
        e.update(dict(zip(key, value)))
    return e



color_dict = {'circle':'FFFFFF', 'ellipse':'FFFFFF',
              'line':'000000', 'path':'000000',
              'rect':'FFFFFF'}

class Draw(object):
    def __init__(self, slide, width, height):
        self.slide = slide
        self.shapes = slide._element.find('.//p:spTree', namespaces=nsmap)
        # TODO: Replace 9999... with slide width and height
        self.x = lambda x: int(float(x) * 9999999 / width)
        self.y = lambda y: int(float(y) * 7777777 / height)

    def _shape_attrs(function):
        def wrapped(self, e):
            par = e.getparent()
            shape = function(self, e)
            tag = function.__name__
            keys = e.keys()


            def styles(keys):
                def clr_grad(color):
                    if color.startswith('rgba('):
                        r, g, b, a = rgba(color)
                        return  '%d' % int((1 - a)*100000 if a < 1 else 100000)
                    elif 'opacity' in keys:
                        return '%d' % int((1 - float(e.get('opacity'))) * 100000)
                    else:
                        return '%d' % 100000


                # TODO: Optimize
                if 'fill' in keys:
                    if e.get('fill') == 'none':
                        shape.spPr.append(a.noFill())
                    else:
                        shape.spPr.append(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('fill')))), 
                            val=str(msclr(e.get('fill'))))))
                elif not 'fill' in keys:
                    if tag not in ['line']:
                        shape.spPr.append(a.solidFill(color(srgbClr='000000')))

                if 'stroke' and 'stroke-width' in keys:
                    shape.spPr.append(a.ln(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('stroke')))),
                        val=str(msclr(e.get('stroke'))))),
                        w=str(int(float(e.get('stroke-width')))*12700/2)))
                    #shape.spPr.append(a.ln(a.solidFill(color(srgbClr=msclr(e.get('stroke')))),
                    #    w=str(int(float(e.get('stroke-width')))*12700/2)))                 
                elif 'stroke' in keys:
                    if e.get('stroke') == 'none':
                        shape.spPr.append(a.ln(a.noFill()))
                    else:
                        shape.spPr.append(a.ln(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('stroke')))), val=str(msclr(e.get('stroke')))))))
                        #shape.spPr.append(a.ln(a.solidFill(color(srgbClr=msclr(e.get('stroke'))))))
                elif 'stroke' and 'fill' not in keys:
                    shape.spPr.append(a.ln(a.solidFill(color(srgbClr='000000'))))
                elif not 'stroke' and 'fill' in keys:
                    if tag in ['rect']:
                        shape.spPr.append(a.ln(a.noFill()))
                    elif tag in ['circle','ellipse']:
                        shape.spPr.append(a.ln(a.solidFill(color(srgbClr=msclr(e.get('fill'))))))
                    elif tag in ['path', 'line']:
                        shape.spPr.append(a.ln(a.solidFill(color(srgbClr='000000'))))
                return shape

            if 'style' in keys:
                e = css_style(e.get('style'))
                keys = e.keys()
                styles(keys)
            #elif par.tag == 'g':
            #    e = par
            #    keys = e.keys()
            #    styles(keys)
            else:
                styles(keys)

            return shape
        return wrapped

    @_shape_attrs
    def circle(self, e):
        x = float(e.get('cx', 0))
        y = float(e.get('cy', 0))
        r = float(e.get('r', 0))
        shp = shape('ellipse', self.x(x - r), self.y(y - r),
                               self.x(2 * r), self.y(2 * r))
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def ellipse(self, e):
        x = float(e.get('cx', 0))
        y = float(e.get('cy', 0))
        rx = float(e.get('rx', 0))
        ry = float(e.get('ry', 0))
        shp = shape('ellipse', self.x(x - rx), self.y(y - ry),
                               self.x(2 * rx), self.y(2 * ry))
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def rect(self, e):
        shp = shape('rect',
            self.x(interpret_str(e.get('x', 0))),
            self.y(interpret_str(e.get('y', 0))),
            self.x(interpret_str(e.get('width', 0))),
            self.y(interpret_str(e.get('height', 0)))
        )
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def line(self, e):
        x1, y1 = self.x(e.get('x1', 0)), self.y(e.get('y1', 0))
        x2, y2 = self.x(e.get('x2', 0)), self.y(e.get('y2', 0))
        shp = shape('line', x1, y1, x2 - x1, y2 - y1)
        self.shapes.append(shp)
        return shp

    def text(self, e):
        keys = e.keys()
        txt = e.text
        def txt_anchor():
            dict = {'hanging':'t', 'middle':'ctr', True:'t', False:'ctr'}
            if 'dominant-baseline' in keys:
                anchor = dict[e.get('dominant-baseline')]
            elif 'dy' in keys:
                em = float(re.findall(".\d+", e.get('dy'))[0]) > 0.5
                anchor = dict[em]
            else:
                anchor = 'ctr'
            return anchor

        def txt_align():
            if 'text-anchor' in keys:
                dict = {'end':'r', 'middle':'ctr', 'start':'l', 'left':'l'}
                align = dict[e.get('text-anchor')]
            else:
                align = 'l'
            return align

        if not e.text:
            return
        shp = shape('rect', self.x(interpret_str(e.get('x', 0))), self.y(interpret_str(e.get('y', 0))), self.x(0), self.y(0))
        def text_style(keys, txt):
            if 'transform' in keys:
                key = e.get('transform')
                shp.append(p.txBody(a.bodyPr(a.normAutofit(fontScale="62500", lnSpcReduction="20000"),
                    a.scene3d(a.camera(a.rot(lat='0', lon='0',
                        rev=str(abs(int(key[(key.find('rotate')+7):-1].split()[0])*60000))),
                        prst='orthographicFront'), a.lightRig(rig='threePt', dir='t')),
                    anchor=txt_anchor(), wrap='none'),
                a.p(a.pPr(algn=txt_align()),
                    a.r(a.t(txt)))))

            elif 'font-size' in keys:
                shp.append(p.txBody(a.bodyPr(anchor='ctr', wrap='none'),
                a.p(a.pPr(algn=txt_align()), a.r(a.rPr(lang='en-US', sz=str(int(float(interpret_str(e.get('font-size')))*100)), dirty='0', smtClean='0'),
                        a.t(txt)))))
            elif 'fill' in keys:
                shp.append(p.txBody(a.bodyPr(a.normAutofit(fontScale="62500", lnSpcReduction="20000"),
                    anchor=txt_anchor(), wrap='none'),
                a.p(a.pPr(algn=txt_align()),
                    a.r(a.rPr( a.solidFill(color(srgbClr=msclr(e.get('fill')))),
                        lang='en-US', dirty='0', smtClean='0'),a.t(txt)))))
            else:
                shp.append(p.txBody(a.bodyPr(a.normAutofit(fontScale="62500", lnSpcReduction="20000"),
                    anchor=txt_anchor(), wrap='none'),
                a.p(a.pPr(algn=txt_align()),
                    a.r(a.t(txt)))))
            return shp

        if 'style' in keys:
            txt = e.text
            e = css_style(e.get('style'))
            keys = e.keys()
            text_style(keys, txt)
        else:
            text_style(keys, txt)

        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def path(self, e):
        pathstr = re_path.findall(e.get('d', ''))
        n, length, cmd, relative, shp = 0, len(pathstr), None, False, None
        x1, y1 = 0, 0
        xy = lambda n: (float(pathstr[n]) + (x1 if relative else 0),
                        float(pathstr[n + 1]) + (y1 if relative else 0))

        shp = cust_shape(x1, y1, 150000, 150000)
        path = a.path(w="150000", h="150000")
        shp.find('.//a:custGeom', namespaces=nsmap).append(
            a.pathLst(path))

        while n < length:
            if pathstr[n].lower() in 'mzlhvcsqta':
                cmd = pathstr[n].lower()
                relative = str.islower(pathstr[n])
                n += 1

            if cmd == 'm':
                x1, y1 = xy(n)
                path.append(a.moveTo(a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 2

            elif cmd == 'z':
                path.append(a.close())

            elif cmd == 'l':
                x1, y1 = xy(n)
                path.append(a.lnTo(a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 2

            elif cmd == 'c':
                xc1, yc1 = xy(n)
                xc2, yc2 = xy(n + 2)
                x1, y1 = xy(n + 4)
                path.append(a.cubicBezTo(a.pt(x=str(self.x(xc1)), y=str(self.y(yc1))),
                    a.pt(x=str(self.x(xc2)), y=str(self.y(yc2))),
                    a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 6

            #TODO blockArc:
            #elif cmd == 'a':
            #    x1, y1 = xy(n)
            #    cx, cy = xy(n + 5)
            #    shp = shape('blockArc', self.x(x1), self.y(y1), self.x(cx), self.y(cy))
            #    n += 7

            elif cmd == 'a':
                wR, hR = xy(n)
                stAng, swAng = xy(n + 5)
                path.append(a.arcTo(
                    wR=str(self.x(wR)), hR=str(self.y(hR)),
                    stAng=str(self.x(stAng)), swAng=str(self.y(swAng))))
                n += 7


        self.shapes.append(shp)
        return shp



def svg2mso(slide, svg, width=940, height=None):
    if width is not None and height is None:
        height = width * 3 / 4
    elif width is None and height is not None:
        width = height * 4 / 3

    # Convert tree into an lxml etree if it's not one
    if not hasattr(svg, 'iter'):
        svg = etree.parse(svg) if hasattr(svg, 'read') else etree.fromstring(svg)

    # Take all the tags and draw it
    draw = Draw(slide, width, height)
    valid_tags = set(tag for tag in dir(draw) if not tag.startswith('_'))
    for e in svg.iter(tag=etree.Element):
        match = re_ns.match(e.tag)
        if not match:
            continue

        tag = match.groups()[-1]
        if tag in valid_tags:
            getattr(draw, tag)(e)


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description=__doc__.strip())
    parser.add_argument('svgfile')
    args = parser.parse_args()

    tree = html.parse(open(args.svgfile))
    #tree = etree.parse(open(args.svgfile))

    from pptx import Presentation

    Presentation = Presentation()
    blank_slidelayout = Presentation.slidelayouts[6]
    slide = Presentation.slides.add_slide(blank_slidelayout)

    svg2mso(slide, tree)
    Presentation.save("test.pptx")
